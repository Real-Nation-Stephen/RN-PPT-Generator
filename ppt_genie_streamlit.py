"""
RN PowerPoint Generator - A Streamlit application for generating PowerPoint presentations from images.

This application helps create PowerPoint presentations by:
- Selecting a folder containing images
- Converting images to PowerPoint slides
- Supporting auto-resize to 16:9 aspect ratio
- Providing a modern, clean interface
- Protected with user authentication

Author: RealNation
Version: 1.0.0
"""

import io
import os
import base64
import streamlit as st
from PIL import Image
from datetime import datetime
import tempfile
import shutil
import glob
from pptx import Presentation
from pptx.util import Inches
import logging
import zipfile
import subprocess
import sys

# Authentication imports
import gspread
from google.auth import exceptions as auth_exceptions
from google.oauth2.service_account import Credentials

# Initialize session state variables
SESSION_STATE_VARS = {
    'selected_folder': None,
    'selected_images': [],
    'generation_complete': False,
    'output_file': None
}

# Initialize all session state variables at once
for var, default_value in SESSION_STATE_VARS.items():
    if var not in st.session_state:
        st.session_state[var] = default_value

# Set up the Streamlit page
st.set_page_config(
    page_title="RN PowerPoint Generator",
    page_icon="assets/App_Icon_PPT.png",
    layout="wide"
)

# Authentication Functions
@st.cache_data(ttl=300)  # Cache for 5 minutes
def load_users_from_sheet():
    """Load users from Google Sheet with caching"""
    try:
        # Create credentials from Streamlit secrets
        creds_dict = dict(st.secrets["service_account"])
        
        # Define the required scopes for Google Sheets API
        scopes = [
            'https://www.googleapis.com/auth/spreadsheets',
            'https://www.googleapis.com/auth/drive'
        ]
        
        credentials = Credentials.from_service_account_info(creds_dict, scopes=scopes)
        
        # Authorize and open the sheet
        gc = gspread.authorize(credentials)
        sheet = gc.open("App Passwords").sheet1
        
        # Get all records
        records = sheet.get_all_records()
        
        # Convert to dictionary
        users = {}
        for record in records:
            name = record.get('Name', '').strip()
            email = record.get('Email', '').strip()
            password = record.get('Password', '').strip()
            image_url = record.get('Image_URL', '').strip()
            
            if name and email and password:
                users[name] = {
                    'email': email,
                    'password': password,
                    'image_url': image_url
                }
        
        return users
    except auth_exceptions.RefreshError:
        st.error("Authentication failed. Please check your Google Service Account credentials.")
        return {}
    except gspread.SpreadsheetNotFound:
        st.error("User database not found. Please check the Google Sheet name.")
        return {}
    except Exception as e:
        st.error(f"Error loading user data: {str(e)}")
        return {}

def convert_google_drive_url(drive_url):
    """Convert Google Drive sharing URL to direct image URL"""
    if not drive_url or "drive.google.com" not in drive_url:
        return drive_url
    
    try:
        if "/file/d/" in drive_url:
            file_id = drive_url.split("/file/d/")[1].split("/")[0]
            return f"https://drive.google.com/uc?export=view&id={file_id}"
        elif "id=" in drive_url:
            file_id = drive_url.split("id=")[1].split("&")[0]
            return f"https://drive.google.com/uc?export=view&id={file_id}"
    except:
        pass
    
    return drive_url

def get_user_gradient(user_name):
    """Get a consistent gradient for a user based on their name"""
    gradients = [
        "#667eea, #764ba2",  # Purple-Blue
        "#f093fb, #f5576c",  # Pink-Red
        "#4facfe, #00f2fe",  # Blue-Cyan
        "#43e97b, #38f9d7",  # Green-Teal
        "#fa709a, #fee140",  # Pink-Yellow
        "#a8edea, #fed6e3",  # Teal-Pink
        "#ffecd2, #fcb69f",  # Peach-Orange
        "#ff9a9e, #fecfef",  # Rose-Purple
    ]
    
    # Use hash of name to consistently assign gradient
    name_hash = hash(user_name.lower()) % len(gradients)
    return gradients[name_hash]

def display_profile_image(image_url, size=100, user_name="User"):
    """Display profile image with fallback"""
    if not image_url:
        # Default avatar with user's initials and unique gradient
        initials = ''.join([name[0].upper() for name in user_name.split()[:2]])
        gradient = get_user_gradient(user_name)
        st.markdown(f"""
        <div style="width: {size}px; height: {size}px; 
                    background: linear-gradient(135deg, {gradient}); 
                    border-radius: 50%; display: flex; align-items: center; 
                    justify-content: center; margin: 0 auto; color: white; 
                    font-size: {size//3}px; font-weight: bold; text-shadow: 0 1px 2px rgba(0,0,0,0.3);">
            {initials}
        </div>
        """, unsafe_allow_html=True)
        return
    
    # Convert Google Drive URLs
    direct_url = convert_google_drive_url(image_url)
    
    try:
        initials = ''.join([name[0].upper() for name in user_name.split()[:2]])
        gradient = get_user_gradient(user_name)
        st.markdown(f"""
        <div style="display: flex; justify-content: center;">
            <img src="{direct_url}" 
                 style="width: {size}px; height: {size}px; border-radius: 50%; 
                        object-fit: cover; border: 3px solid #2b1e66;" 
                 onerror="this.style.display='none'; this.nextElementSibling.style.display='flex';"
                 alt="{user_name}">
            <div style="width: {size}px; height: {size}px; 
                        background: linear-gradient(135deg, {gradient}); 
                        border-radius: 50%; display: none; align-items: center; 
                        justify-content: center; color: white; font-size: {size//3}px; 
                        font-weight: bold; text-shadow: 0 1px 2px rgba(0,0,0,0.3);">
                {initials}
            </div>
        </div>
        """, unsafe_allow_html=True)
    except:
        # Fallback to initials with gradient
        initials = ''.join([name[0].upper() for name in user_name.split()[:2]])
        gradient = get_user_gradient(user_name)
        st.markdown(f"""
        <div style="width: {size}px; height: {size}px; 
                    background: linear-gradient(135deg, {gradient}); 
                    border-radius: 50%; display: flex; align-items: center; 
                    justify-content: center; margin: 0 auto; color: white; 
                    font-size: {size//3}px; font-weight: bold; text-shadow: 0 1px 2px rgba(0,0,0,0.3);">
            {initials}
        </div>
        """, unsafe_allow_html=True)

def check_authentication():
    """Check if user is authenticated"""
    try:
        if 'authenticated' not in st.session_state:
            st.session_state.authenticated = False
            st.session_state.selected_user = None
        
        if not st.session_state.authenticated:
            return show_login()
        return True
    except Exception as e:
        st.error(f"Authentication error: {str(e)}")
        return False

def show_login():
    """Display the polished 3-panel login screen"""
    users = load_users_from_sheet()
    if not users:
        st.error("Unable to load user data. Please contact support.")
        return False
    
    # Initialize selected user
    if 'selected_user' not in st.session_state or st.session_state.selected_user is None:
        st.session_state.selected_user = list(users.keys())[0] if users else None
    
    selected_user = st.session_state.selected_user
    
    # Safety check - if selected_user is still None or not in users, reset it
    if selected_user is None or selected_user not in users:
        if users:
            selected_user = list(users.keys())[0]
            st.session_state.selected_user = selected_user
        else:
            st.error("No users found in the system.")
            return False

    # Custom CSS for sign-in button
    st.markdown("""
    <style>
    .stButton > button[kind="primary"] {
        background-color: #2b1e66 !important;
        border-color: #2b1e66 !important;
        color: white !important;
    }
    .stButton > button[kind="primary"]:hover {
        background-color: #1e1547 !important;
        border-color: #1e1547 !important;
        color: white !important;
    }
    .stButton > button[kind="primary"]:focus {
        background-color: #1e1547 !important;
        border-color: #1e1547 !important;
        color: white !important;
        box-shadow: 0 0 0 2px rgba(43, 30, 102, 0.3) !important;
    }
    </style>
    """, unsafe_allow_html=True)

    st.markdown("<br><br><br>", unsafe_allow_html=True)  # Top spacing

    # Container to constrain width
    outer_col1, main_container, outer_col3 = st.columns([2, 4, 2])
    
    with main_container:
        # TOP PANEL - 3 columns with dividers
        col1, div1, col2, div2, col3 = st.columns([1, 0.1, 2, 0.1, 1])
        
        # Column 1: Profile Image
        with col1:
            if selected_user in users:
                user_image_url = users[selected_user].get('image_url', '')
                display_profile_image(user_image_url, size=100, user_name=selected_user)
        
        # Divider 1
        with div1:
            st.markdown('<div style="width: 1px; height: 120px; background-color: #ddd; margin: 0 auto;"></div>', unsafe_allow_html=True)
        
        # Column 2: App Title
        with col2:
            st.markdown("### RN PowerPoint Generator")
            st.caption("Welcome back! Please sign in to continue.")
        
        # Divider 2  
        with div2:
            st.markdown('<div style="width: 1px; height: 120px; background-color: #ddd; margin: 0 auto;"></div>', unsafe_allow_html=True)
        
        # Column 3: Account Info
        with col3:
            if selected_user in users:
                st.markdown(f"**{selected_user}**")
                st.caption(f"📧 {users[selected_user].get('email', 'No email')}")
    
        st.markdown("<br>", unsafe_allow_html=True)
        
        # MIDDLE PANEL - User Selection
        st.markdown("---")
        st.markdown("#### 👤 Select User")
        user_list = list(users.keys())
        try:
            current_index = user_list.index(selected_user) if selected_user in user_list else 0
        except (ValueError, IndexError):
            current_index = 0
            
        new_selection = st.selectbox(
            "Choose your name:",
            options=user_list,
            index=current_index,
            key="user_dropdown",
            label_visibility="collapsed"
        )
        
        if new_selection != st.session_state.selected_user:
            st.session_state.selected_user = new_selection
            st.rerun()
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # BOTTOM PANEL - Password Entry
        st.markdown("---")
        with st.form("login_form"):
            st.markdown("#### 🔒 Enter Password")
            password = st.text_input(
                "Password:", 
                type="password", 
                placeholder="Enter your password...",
                label_visibility="collapsed"
            )
            
            st.markdown("<br>", unsafe_allow_html=True)
            
            login_button = st.form_submit_button(
                "🚀 Sign In", 
                type="primary", 
                use_container_width=True
            )
            
            if login_button:
                current_user = st.session_state.get('selected_user', list(users.keys())[0])
                if current_user in users and password == users[current_user]['password']:
                    st.session_state.authenticated = True
                    st.session_state.current_user = current_user
                    st.session_state.user_email = users[current_user]['email']
                    st.session_state.user_image = users[current_user].get('image_url', '')
                    st.success(f"🎉 Welcome back, {current_user}!")
                    st.rerun()
                else:
                    st.error("❌ Invalid credentials. Please try again.")
        
        # Footer
        st.markdown("""
        <div style="text-align: center; margin-top: 2rem; color: #888; font-size: 0.9rem;">
            Need help? Contact your administrator.
        </div>
        """, unsafe_allow_html=True)
    
    return False

def show_logout_button():
    """Show polished sidebar profile with logout"""
    with st.sidebar:
        st.markdown("---")
        
        # Profile image centered
        user_image = st.session_state.get('user_image', '')
        current_user = st.session_state.get('current_user', '')
        
        # Center the profile image
        profile_col1, profile_col2, profile_col3 = st.columns([1, 1, 1])
        with profile_col2:
            display_profile_image(user_image, size=80, user_name=current_user)
        
        # User details centered with white pills
        st.markdown(f"""
        <div style="text-align: center; margin: 1rem 0;">
            <div style="background-color: white; border-radius: 15px; padding: 0.3rem 0.8rem; margin: 0.3rem auto; display: inline-block; color: #333; font-weight: 500; font-size: 0.9rem; box-shadow: 0 1px 3px rgba(0,0,0,0.1);">
                {st.session_state.get('current_user', 'Unknown')}
            </div>
            <br>
            <div style="background-color: white; border-radius: 15px; padding: 0.3rem 0.8rem; margin: 0.3rem auto; display: inline-block; color: #666; font-size: 0.8rem; box-shadow: 0 1px 3px rgba(0,0,0,0.1);">
                📧 {st.session_state.get('user_email', 'Unknown')}
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        if st.button("🚪 Logout", type="secondary", use_container_width=True):
            # Clear authentication
            st.session_state.authenticated = False
            if 'current_user' in st.session_state:
                del st.session_state.current_user
            if 'user_email' in st.session_state:
                del st.session_state.user_email
            if 'user_image' in st.session_state:
                del st.session_state.user_image
            st.rerun()
        
        st.markdown("---")

def show_user_info_block():
    """Show user info block for main page display"""
    if not st.session_state.get('authenticated', False):
        return
    
    current_user = st.session_state.get('current_user', 'Unknown')
    user_email = st.session_state.get('user_email', 'Unknown')
    user_image = st.session_state.get('user_image', '')
    
    # Create a nice info block
    col1, col2 = st.columns([1, 4])
    
    with col1:
        display_profile_image(user_image, size=60, user_name=current_user)
    
    with col2:
        st.markdown(f"""
        <div style="padding: 0.5rem 0;">
            <div style="background-color: white; border-radius: 12px; padding: 0.2rem 0.6rem; margin: 0.2rem 0; display: inline-block; color: #333; font-weight: 500; font-size: 0.9rem; box-shadow: 0 1px 3px rgba(0,0,0,0.1);">
                👤 {current_user}
            </div>
            <br>
            <div style="background-color: white; border-radius: 12px; padding: 0.2rem 0.6rem; margin: 0.2rem 0; display: inline-block; color: #666; font-size: 0.8rem; box-shadow: 0 1px 3px rgba(0,0,0,0.1);">
                📧 {user_email}
            </div>
        </div>
        """, unsafe_allow_html=True)

# Authentication check - MUST be early, before any other st.* calls
if not check_authentication():
    st.stop()

# Show logout button in sidebar
show_logout_button()

# Global padding style (matching contrast checker)
st.markdown("""
    <style>
        .main .block-container {
            padding-right: 5rem;
        }
        
        /* Custom button styling */
        .stButton > button {
            background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
            color: white;
            border: none;
            border-radius: 8px;
            padding: 0.5rem 1rem;
            font-weight: 500;
            transition: all 0.3s ease;
        }
        
        .stButton > button:hover {
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(102, 126, 234, 0.4);
        }
        
        /* Special styling for generate button */
        div[data-testid="stButton"] button[kind="primary"] {
            background-color: #2b1e66 !important;
            color: #FFFFFF !important;
            border: none !important;
        }
        
        div[data-testid="stButton"] button[kind="primary"]:hover {
            background-color: #1e1448 !important;
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(43, 30, 102, 0.4);
        }
        
        /* Toggle switch styling */
        .stCheckbox > label {
            font-size: 1.1rem;
            font-weight: 500;
        }
        
        /* Success/Error message styling */
        .success-message {
            padding: 1rem;
            background: linear-gradient(90deg, #56ab2f 0%, #a8e6cf 100%);
            color: white;
            border-radius: 8px;
            margin: 1rem 0;
            text-align: center;
            font-weight: 500;
        }
        
        .error-message {
            padding: 1rem;
            background: linear-gradient(90deg, #ff416c 0%, #ff4b2b 100%);
            color: white;
            border-radius: 8px;
            margin: 1rem 0;
            text-align: center;
            font-weight: 500;
        }
        
        /* File upload area styling */
        .uploadedFile {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            border-radius: 8px;
            padding: 1rem;
            margin: 0.5rem 0;
        }
    </style>
""", unsafe_allow_html=True)

# Add sidebar with logo (matching contrast checker)
st.sidebar.markdown("""
    <style>
        /* Container for the logo */
        .sidebar-logo {
            position: fixed;
            bottom: 0;
            padding: 2rem 0.5rem 2rem 0.5rem;
            z-index: 1000;
            width: calc(100% - 1rem);
        }
        
        /* Logo image styling */
        .sidebar-logo img {
            width: calc(100% - 1rem);
            max-width: 200px;
            height: auto;
            display: block;
        }

        /* Theme-aware logo display */
        @media (prefers-color-scheme: light) {
            .sidebar-logo .logo-light { display: block !important; }
            .sidebar-logo .logo-dark { display: none !important; }
        }
        @media (prefers-color-scheme: dark) {
            .sidebar-logo .logo-light { display: none !important; }
            .sidebar-logo .logo-dark { display: block !important; }
        }
        
        /* Default fallback - show light logo */
        .sidebar-logo .logo-light { display: block; }
        .sidebar-logo .logo-dark { display: none; }
    </style>
""", unsafe_allow_html=True)

# Logo handling
try:
    # Try to load actual logo files
    if os.path.exists("assets/light-mode-logo.png"):
        with open("assets/light-mode-logo.png", "rb") as f:
            logo_light_base64 = base64.b64encode(f.read()).decode()
    else:
        # Fallback to a simple SVG logo
        logo_light_svg = """
        <svg width="200" height="80" viewBox="0 0 200 80" fill="none" xmlns="http://www.w3.org/2000/svg">
            <rect width="200" height="80" rx="10" fill="#2E3440"/>
            <text x="100" y="45" text-anchor="middle" fill="#D8DEE9" font-family="Arial" font-size="18" font-weight="bold">PPT Generator</text>
        </svg>
        """
        logo_light_base64 = base64.b64encode(logo_light_svg.encode()).decode()

    if os.path.exists("assets/dark-mode-logo.png"):
        with open("assets/dark-mode-logo.png", "rb") as f:
            logo_dark_base64 = base64.b64encode(f.read()).decode()
    else:
        # Fallback dark logo
        logo_dark_svg = """
        <svg width="200" height="80" viewBox="0 0 200 80" fill="none" xmlns="http://www.w3.org/2000/svg">
            <rect width="200" height="80" rx="10" fill="url(#gradient)"/>
            <text x="100" y="45" text-anchor="middle" fill="white" font-family="Arial" font-size="18" font-weight="bold">PPT Generator</text>
            <defs>
                <linearGradient id="gradient" x1="0%" y1="0%" x2="100%" y2="0%">
                    <stop offset="0%" style="stop-color:#667eea"/>
                    <stop offset="100%" style="stop-color:#764ba2"/>
                </linearGradient>
            </defs>
        </svg>
        """
        logo_dark_base64 = base64.b64encode(logo_dark_svg.encode()).decode()

    # Display logos with theme awareness
    st.sidebar.markdown(f"""
        <div class="sidebar-logo">
            <img src="data:image/png;base64,{logo_light_base64}" class="logo-light">
            <img src="data:image/png;base64,{logo_dark_base64}" class="logo-dark">
        </div>
    """, unsafe_allow_html=True)
    
except Exception as e:
    st.sidebar.write("Logo loading error")

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def get_default_slide_size():
    """Get default 16:9 slide size (1920x1080)"""
    width_inches = 13.33  # 1920/144
    height_inches = 7.5   # 1080/144
    return width_inches, height_inches

def check_image_file(image_path):
    """Check if image file is valid and can be opened"""
    try:
        with Image.open(image_path) as img:
            img.verify()
        return True
    except Exception as e:
        logger.warning(f"Invalid or corrupted image file {image_path}: {str(e)}")
        return False

def generate_ppt_from_images(image_files, auto_resize=True):
    """Generate PowerPoint presentation from uploaded images"""
    try:
        # Create presentation
        prs = Presentation()
        
        # Set default 16:9 size
        width_inches, height_inches = get_default_slide_size()
        prs.slide_width = Inches(width_inches)
        prs.slide_height = Inches(height_inches)
        
        # Process each image
        for i, uploaded_file in enumerate(image_files):
            try:
                # Create a temporary file for the image
                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_file.name.split('.')[-1]}") as tmp_file:
                    tmp_file.write(uploaded_file.getvalue())
                    tmp_file_path = tmp_file.name
                
                # Add a slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
                
                # Remove any default shapes
                for shape in list(slide.shapes):
                    slide.shapes._spTree.remove(shape._element)
                
                if auto_resize:
                    # Add picture fitting to slide size (16:9)
                    slide.shapes.add_picture(
                        tmp_file_path,
                        0, 0,
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                else:
                    # Add picture with original dimensions
                    slide.shapes.add_picture(tmp_file_path, 0, 0)
                
                # Clean up temporary file
                os.unlink(tmp_file_path)
                
            except Exception as e:
                logger.error(f"Error processing image {uploaded_file.name}: {str(e)}")
                continue
        
        # Save to memory buffer
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error generating presentation: {str(e)}")
        raise

# Add title and description
st.title("RN PowerPoint Generator")
st.markdown("""
This tool creates PowerPoint presentations from your images. 
Upload images and generate a professional presentation with automatic 16:9 formatting.
""")

# Optional: Show current user info
with st.expander("👤 Current User", expanded=False):
    show_user_info_block()

# File uploader for images
uploaded_files = st.file_uploader(
    "Upload Images", 
    type=['jpg', 'jpeg', 'png', 'JPG', 'JPEG', 'PNG'],
    accept_multiple_files=True,
    help="Select multiple image files to include in your presentation"
)

# Options section
st.subheader("Options")
col1, col2 = st.columns(2)

with col1:
    auto_resize = st.checkbox(
        "Auto Resize to 16:9", 
        value=True, 
        help="Automatically resize images to fit 16:9 aspect ratio"
    )

with col2:
    if uploaded_files:
        st.metric("Images Selected", len(uploaded_files))

# Preview section
if uploaded_files:
    st.subheader("Image Preview")
    
    # Show preview of first few images
    preview_cols = st.columns(min(4, len(uploaded_files)))
    for i, uploaded_file in enumerate(uploaded_files[:4]):
        with preview_cols[i]:
            try:
                image = Image.open(uploaded_file)
                st.image(image, caption=uploaded_file.name, use_container_width=True)
            except Exception as e:
                st.error(f"Error loading {uploaded_file.name}")
    
    if len(uploaded_files) > 4:
        st.info(f"... and {len(uploaded_files) - 4} more images")

# Generation section
st.subheader("Generate PowerPoint")

if st.button("🚀 Generate Presentation", disabled=not uploaded_files, type="primary"):
    if not uploaded_files:
        st.error("Please upload at least one image file.")
    else:
        try:
            with st.spinner("Generating PowerPoint presentation..."):
                # Generate the presentation
                ppt_buffer = generate_ppt_from_images(uploaded_files, auto_resize)
                
                # Success message
                st.markdown("""
                    <div class="success-message">
                        ✅ PowerPoint presentation generated successfully!
                    </div>
                """, unsafe_allow_html=True)
                
                # Download button
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"presentation_{timestamp}.pptx"
                
                st.download_button(
                    label="📥 Download PowerPoint",
                    data=ppt_buffer.getvalue(),
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                
                # Show summary
                st.info(f"Created presentation with {len(uploaded_files)} slides")
                
        except Exception as e:
            st.markdown(f"""
                <div class="error-message">
                    ❌ Error generating presentation: {str(e)}
                </div>
            """, unsafe_allow_html=True)

# Instructions section
with st.expander("📖 How to Use", expanded=False):
    st.markdown("""
    ### Step-by-step Guide:
    
    1. **Upload Images**: Click "Upload Images" and select multiple image files (JPG, PNG)
    2. **Configure Options**: 
       - ✅ **Auto Resize to 16:9**: Automatically formats images for widescreen presentations
       - ❌ **Auto Resize to 16:9**: Keeps original image dimensions
    3. **Preview**: Check the preview to ensure your images are loaded correctly
    4. **Generate**: Click "🚀 Generate Presentation" to create your PowerPoint
    5. **Download**: Use the download button to save your presentation
    
    ### Tips:
    NB: Ensure when exporting to use the format 'file name' with an underscore at the end. Output images should thus be labelled "file name_", "file name_1", "file name_2" etc. This is vital to ensure correct sequencing.
    
    NB: The folder of images MUST be saved locally on your device (usually desktop, download, documents) - NOT saved in a Google Drive location. Just copy and paste it onto your desktop or downloads and delete once the presentation has been generated.
    """)

# Footer
st.markdown("---")
st.markdown(
    f"**RN PowerPoint Generator** | Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
) 